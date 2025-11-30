#!/usr/bin/env python3
"""
Generate presentation slides for Type-Theoretic Ontology System
Using python-pptx library with Classic Blue + Teal color palette
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Color Palette: Classic Blue + Teal
COLOR_NAVY = RGBColor(28, 40, 51)        # #1C2833
COLOR_SLATE = RGBColor(46, 64, 83)       # #2E4053
COLOR_TEAL = RGBColor(94, 168, 167)      # #5EA8A7
COLOR_SILVER = RGBColor(170, 183, 184)   # #AAB7B8
COLOR_LIGHT = RGBColor(244, 246, 246)    # #F4F6F6
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_CODE_BG = RGBColor(236, 240, 241)  # #ECF0F1

def create_presentation():
    """Create the presentation with all slides"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    # Slide 1: Title
    slide1 = add_title_slide(prs)

    # Slide 2: Overview
    slide2 = add_overview_slide(prs)

    # Slide 3: Section 1 - Theory
    slide3 = add_section_slide(prs, "第1部：理論的基礎", "Span/Cospan理論と型理論")

    # Slide 4: Span/Cospan Theory
    slide4 = add_span_cospan_slide(prs)

    # Slide 5: Type Theory Basics
    slide5 = add_type_theory_slide(prs)

    # Slide 6: Type Inhabitation
    slide6 = add_type_inhabitation_slide(prs)

    # Slide 7: Automatic Discovery
    slide7 = add_auto_discovery_slide(prs)

    # Slide 8: Section 2 - YAML
    slide8 = add_section_slide(prs, "第2部：YAMLベースの実装", "簡易DSLによる探索システム")

    # Slide 9: YAML Catalog
    slide9 = add_yaml_catalog_slide(prs)

    # Slide 10: Search Algorithm
    slide10 = add_search_algorithm_slide(prs)

    # Slide 11: Validation Results
    slide11 = add_validation_results_slide(prs)

    # Slide 12: Section 3 - DSL Design
    slide12 = add_section_slide(prs, "第3部：本格的なDSL設計", "拡張可能なDSLと実行エンジン")

    # Slide 13: DSL Syntax
    slide13 = add_dsl_syntax_slide(prs)

    # Slide 14: Implementation Types
    slide14 = add_implementation_types_slide(prs)

    # Slide 15: Execution Engine
    slide15 = add_execution_engine_slide(prs)

    # Slide 16: Provenance
    slide16 = add_provenance_slide(prs)

    # Slide 17: Section 4 - CFP
    slide17 = add_section_slide(prs, "第4部：CFP計算の実例", "Carbon Footprint計算での適用")

    # Slide 18: CFP Problem
    slide18 = add_cfp_problem_slide(prs)

    # Slide 19: Search Process
    slide19 = add_search_process_slide(prs)

    # Slide 20: GHG Scope 123
    slide20 = add_ghg_scope_slide(prs)

    # Slide 21: Analysis Results
    slide21 = add_analysis_results_slide(prs)

    # Slide 22: Metrics Analysis
    slide22 = add_metrics_analysis_slide(prs)

    # Slide 23: Section 5 - Summary
    slide23 = add_section_slide(prs, "第5部：まとめと今後の展望", "")

    # Slide 24: Achievements
    slide24 = add_achievements_slide(prs)

    # Slide 25: Future Work
    slide25 = add_future_work_slide(prs)

    # Slide 26: Conclusion
    slide26 = add_conclusion_slide(prs)

    return prs

def add_title_slide(prs):
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Navy background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "型理論ベース\nオントロジー探索・合成システム"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_LIGHT
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Type-Theoretic Ontology Exploration and Synthesis"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(20)
    subtitle_p.font.color.rgb = COLOR_TEAL
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = "カテゴリ理論 × 型理論 × オントロジー工学"
    footer_p = footer_frame.paragraphs[0]
    footer_p.font.size = Pt(18)
    footer_p.font.color.rgb = COLOR_LIGHT
    footer_p.alignment = PP_ALIGN.CENTER

    return slide

def add_overview_slide(prs):
    """Slide 2: Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    # Title
    add_title(slide, "プロジェクト概要")

    # Content
    content_top = Inches(1.2)

    # Purpose section
    add_heading(slide, "目的", Inches(0.5), content_top, COLOR_TEAL)
    add_text(slide, "異なるオントロジー間の変換を、型理論とカテゴリ理論を用いて自動的に発見・合成する",
             Inches(0.5), content_top + Inches(0.3), Inches(9), Inches(0.4))

    # Components section
    add_heading(slide, "主要コンポーネント", Inches(0.5), content_top + Inches(0.9), COLOR_TEAL)
    bullets = [
        "Span/Cospan理論に基づくオントロジーアライメント",
        "型充足問題による関数合成の自動発見",
        "DSLによる宣言的なカタログ定義",
        "実行エンジン（SPARQL/REST/Formula対応）",
        "PROV-O準拠のProvenance生成"
    ]
    add_bullet_list(slide, bullets, Inches(0.5), content_top + Inches(1.2), Inches(9), Inches(2.5))

    return slide

def add_section_slide(prs, title, subtitle):
    """Section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Slate background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_SLATE

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_LIGHT
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.font.size = Pt(20)
        subtitle_p.font.color.rgb = COLOR_TEAL
        subtitle_p.alignment = PP_ALIGN.CENTER

    return slide

def add_span_cospan_slide(prs):
    """Slide 4: Span/Cospan Theory"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "Span/Cospan理論")

    top = Inches(1.2)

    add_heading(slide, "Span構造", Inches(0.5), top, COLOR_TEAL)
    add_text(slide, "2つのオントロジーを中間オントロジーで接続",
             Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.3))
    add_code(slide, "A ← M → B", Inches(0.5), top + Inches(0.6), Inches(9), Inches(0.4))

    bullets1 = [
        "A, B: ソース・ターゲットオントロジー",
        "M: 中間オントロジー（マッピング仲介者）",
        "射: オントロジー準同型（構造を保つ写像）"
    ]
    add_bullet_list(slide, bullets1, Inches(0.5), top + Inches(1.1), Inches(9), Inches(0.9))

    add_heading(slide, "Cospan構造", Inches(0.5), top + Inches(2.2), COLOR_TEAL)
    add_text(slide, "2つのオントロジーから共通の統合オントロジーへ",
             Inches(0.5), top + Inches(2.5), Inches(9), Inches(0.3))
    add_code(slide, "A → U ← B", Inches(0.5), top + Inches(2.8), Inches(9), Inches(0.4))

    add_text(slide, "Pullback操作により、最も精密な整合性を保つマッピングを構築",
             Inches(0.5), top + Inches(3.3), Inches(9), Inches(0.3))

    return slide

def add_type_theory_slide(prs):
    """Slide 5: Type Theory Basics"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "型理論の基礎")

    top = Inches(1.2)

    add_heading(slide, "型システム", Inches(0.5), top, COLOR_TEAL)
    bullets1 = [
        "型 (Type): データの分類・制約を表現",
        "項 (Term): 型に属する具体的な値",
        "判断 (Judgment): t : T （項tが型Tに属する）"
    ]
    add_bullet_list(slide, bullets1, Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.9))

    add_heading(slide, "関数型", Inches(0.5), top + Inches(1.4), COLOR_TEAL)
    add_text(slide, "関数は型から型への写像",
             Inches(0.5), top + Inches(1.7), Inches(9), Inches(0.3))
    add_code(slide, "f : A → B", Inches(0.5), top + Inches(2.0), Inches(9), Inches(0.4))
    add_text(slide, "関数fは型Aの値を受け取り、型Bの値を返す",
             Inches(0.5), top + Inches(2.5), Inches(9), Inches(0.3))

    add_heading(slide, "関数合成", Inches(0.5), top + Inches(3.0), COLOR_TEAL)
    add_text(slide, "f : A → B と g : B → C があるとき",
             Inches(0.5), top + Inches(3.3), Inches(9), Inches(0.3))
    add_code(slide, "g ∘ f : A → C", Inches(0.5), top + Inches(3.6), Inches(9), Inches(0.4))

    return slide

def add_type_inhabitation_slide(prs):
    """Slide 6: Type Inhabitation Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "型充足問題")

    top = Inches(1.2)

    add_heading(slide, "問題定義", Inches(0.5), top, COLOR_TEAL)
    add_text(slide, "与えられた型Tに対して、その型を持つ項を見つける問題",
             Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.3))
    add_code(slide, "Find t such that t : T", Inches(0.5), top + Inches(0.6), Inches(9), Inches(0.4))

    add_heading(slide, "本システムへの適用", Inches(0.5), top + Inches(1.2), COLOR_TEAL)
    bullets1 = [
        "Goal: ソース型 S からターゲット型 T への変換",
        "Catalog: 利用可能な関数の集合",
        "Task: S → T という型を満たす関数合成を発見"
    ]
    add_bullet_list(slide, bullets1, Inches(0.5), top + Inches(1.5), Inches(9), Inches(0.9))

    add_heading(slide, "探索アルゴリズム", Inches(0.5), top + Inches(2.6), COLOR_TEAL)
    add_text(slide, "Dijkstra法ベースの後方探索",
             Inches(0.5), top + Inches(2.9), Inches(9), Inches(0.3))
    bullets2 = [
        "ゴール型Tから開始し、Sに到達するパスを探索",
        "コスト最小・信頼度最大のパスを優先"
    ]
    add_bullet_list(slide, bullets2, Inches(0.5), top + Inches(3.2), Inches(9), Inches(0.6))

    return slide

def add_auto_discovery_slide(prs):
    """Slide 7: Automatic Discovery"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "中間型の自動発見")

    top = Inches(1.2)

    add_heading(slide, "従来の課題", Inches(0.5), top, COLOR_TEAL)
    add_text(slide, "オントロジーマッピングでは、中間オントロジーを手動で設計する必要があった",
             Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.4))

    add_heading(slide, "本システムのアプローチ", Inches(0.5), top + Inches(0.9), COLOR_TEAL)
    bullets1 = [
        "型充足問題として定式化",
        "カタログから関数を選択し、自動的に合成",
        "複数の中間型を経由するパスも発見可能"
    ]
    add_bullet_list(slide, bullets1, Inches(0.5), top + Inches(1.2), Inches(9), Inches(0.9))

    add_heading(slide, "例", Inches(0.5), top + Inches(2.3), COLOR_TEAL)
    code_text = """Goal: PurchaseAmount -> kg-CO2

Path found:
PurchaseAmount -> MonetaryValue -> EnergyConsumption -> kg-CO2

Functions:
  purchaseToMonetary : PurchaseAmount -> MonetaryValue
  monetaryToEnergy   : MonetaryValue -> EnergyConsumption
  energyToEmissions  : EnergyConsumption -> kg-CO2"""
    add_code(slide, code_text, Inches(0.5), top + Inches(2.6), Inches(9), Inches(1.5))

    return slide

def add_yaml_catalog_slide(prs):
    """Slide 9: YAML Catalog"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "YAMLカタログ定義")

    top = Inches(1.2)

    add_heading(slide, "構造", Inches(0.5), top, COLOR_TEAL)
    code_text = """types:
  - name: MonetaryValue
    unit: JPY
  - name: EnergyConsumption
    unit: kWh

functions:
  - name: purchaseToEnergy
    source: MonetaryValue
    target: EnergyConsumption
    cost: 0.5
    confidence: 0.9
    implementation:
      type: formula
      expr: "x * 0.002\""""
    add_code(slide, code_text, Inches(0.5), top + Inches(0.3), Inches(9), Inches(2))

    add_heading(slide, "特徴", Inches(0.5), top + Inches(2.5), COLOR_TEAL)
    bullets = [
        "人間が読みやすい宣言的記述",
        "型、関数、コスト、信頼度を明示",
        "実装方法（formula/SPARQL/REST）を指定可能"
    ]
    add_bullet_list(slide, bullets, Inches(0.5), top + Inches(2.8), Inches(9), Inches(0.9))

    return slide

def add_search_algorithm_slide(prs):
    """Slide 10: Search Algorithm"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "探索アルゴリズム")

    top = Inches(1.2)

    add_heading(slide, "Dijkstra後方探索", Inches(0.5), top, COLOR_TEAL)
    code_text = """def synthesize_backward(catalog, source_type, goal_type):
    pq = [(0.0, 0.0, counter, goal_type, [])]
    visited = set()

    while pq:
        cum_cost, cum_conf, _, current_type, path = heappop(pq)

        if current_type == source_type:
            return path  # Solution found

        for func in catalog.functions:
            if func.target == current_type:
                next_type = func.source
                new_cost = cum_cost + func.cost
                new_conf = cum_conf + (1.0 - func.confidence)
                heappush(pq, (new_cost, new_conf, counter,
                         next_type, [func] + path))"""
    add_code(slide, code_text, Inches(0.5), top + Inches(0.3), Inches(9), Inches(2.2))

    bullets = [
        "優先度キュー: (コスト, 信頼度損失, カウンタ, 型, パス)",
        "ゴールから開始、ソースに到達するまで後方探索"
    ]
    add_bullet_list(slide, bullets, Inches(0.5), top + Inches(2.7), Inches(9), Inches(0.6))

    return slide

def add_validation_results_slide(prs):
    """Slide 11: Validation Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "検証結果")

    top = Inches(1.2)

    add_heading(slide, "テストシナリオ", Inches(0.5), top, COLOR_TEAL)
    bullets1 = [
        "CFP計算：購買金額 → CO2排出量",
        "エネルギー換算：電力消費 → ガス消費",
        "複数中間型を経由するパス"
    ]
    add_bullet_list(slide, bullets1, Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.9))

    add_heading(slide, "成果", Inches(0.5), top + Inches(1.4), COLOR_TEAL)
    bullets2 = [
        "100% のパス発見成功率",
        "平均 2-3ホップ で目標型に到達",
        "コスト最小パスの正確な選択",
        "信頼度を考慮した最適化"
    ]
    add_bullet_list(slide, bullets2, Inches(0.5), top + Inches(1.7), Inches(9), Inches(1.2))

    add_heading(slide, "限界", Inches(0.5), top + Inches(3.1), COLOR_TEAL)
    bullets3 = [
        "集約関数（複数入力）の扱いが単純化されている",
        "単位変換が型システムに統合されていない"
    ]
    add_bullet_list(slide, bullets3, Inches(0.5), top + Inches(3.4), Inches(9), Inches(0.6))

    return slide

def add_dsl_syntax_slide(prs):
    """Slide 13: DSL Syntax"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "DSL構文")

    top = Inches(1.2)

    add_heading(slide, "型定義", Inches(0.5), top, COLOR_TEAL)
    code1 = """type EmissionFactor [unit=kg-CO2/kWh] {
  doc: "Emission factor for energy consumption"
  metadata: {
    domain: "GHG Protocol"
    version: "2024"
  }
}"""
    add_code(slide, code1, Inches(0.5), top + Inches(0.3), Inches(9), Inches(1.1))

    add_heading(slide, "関数定義", Inches(0.5), top + Inches(1.6), COLOR_TEAL)
    code2 = """fn energyToEmissions {
  sig: EnergyConsumption -> EmissionFactor -> kg-CO2
  impl: formula("energy * factor")
  cost: 0.1
  confidence: 0.95
  provenance: {
    method: "IPCC Guidelines"
    source: "IEA Database"
  }
}"""
    add_code(slide, code2, Inches(0.5), top + Inches(1.9), Inches(9), Inches(1.5))

    return slide

def add_implementation_types_slide(prs):
    """Slide 14: Implementation Types"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "実装の種類")

    top = Inches(1.2)

    add_heading(slide, "1. Formula実装", Inches(0.5), top, COLOR_TEAL)
    add_text(slide, "数式による直接計算", Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.3))
    add_code(slide, 'impl: formula("x * 0.5 + 100")', Inches(0.5), top + Inches(0.6), Inches(9), Inches(0.4))

    add_heading(slide, "2. SPARQL実装", Inches(0.5), top + Inches(1.2), COLOR_TEAL)
    add_text(slide, "RDFデータベースへのクエリ", Inches(0.5), top + Inches(1.5), Inches(9), Inches(0.3))
    code2 = """impl: sparql {
  endpoint: "http://example.org/sparql"
  query: "SELECT ?value WHERE { ?s ex:property ?value }"
}"""
    add_code(slide, code2, Inches(0.5), top + Inches(1.8), Inches(9), Inches(0.7))

    add_heading(slide, "3. REST実装", Inches(0.5), top + Inches(2.7), COLOR_TEAL)
    add_text(slide, "外部APIの呼び出し", Inches(0.5), top + Inches(3.0), Inches(9), Inches(0.3))
    code3 = """impl: rest {
  url: "https://api.example.com/convert"
  method: "POST"
  body_template: '{"value": {{input}}}'
}"""
    add_code(slide, code3, Inches(0.5), top + Inches(3.3), Inches(9), Inches(0.7))

    return slide

def add_execution_engine_slide(prs):
    """Slide 15: Execution Engine"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "実行エンジン")

    top = Inches(1.2)

    add_heading(slide, "機能", Inches(0.5), top, COLOR_TEAL)
    bullets = [
        "Executor: 関数パスの実行（formula/SPARQL/REST対応）",
        "UnitConverter: 単位変換の自動実行",
        "ProvenanceGenerator: PROV-O準拠のProvenance生成"
    ]
    add_bullet_list(slide, bullets, Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.9))

    add_heading(slide, "実行例", Inches(0.5), top + Inches(1.4), COLOR_TEAL)
    code = """result = executor.execute(path, input_value)
# result = {
#   "value": 123.45,
#   "unit": "kg-CO2",
#   "provenance": {
#     "activity": "prov:Activity_001",
#     "wasGeneratedBy": [...],
#     "used": [...]
#   }
# }"""
    add_code(slide, code, Inches(0.5), top + Inches(1.7), Inches(9), Inches(1.5))

    add_text(slide, "実行の各ステップでProvenanceを記録し、トレーサビリティを確保",
             Inches(0.5), top + Inches(3.3), Inches(9), Inches(0.4))

    return slide

def add_provenance_slide(prs):
    """Slide 16: Provenance"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "PROV-O Provenance")

    top = Inches(1.2)

    add_heading(slide, "W3C PROV-O準拠", Inches(0.5), top, COLOR_TEAL)
    bullets = [
        "Entity: 入力・出力データ",
        "Activity: 関数の実行",
        "wasGeneratedBy: データ生成の因果関係",
        "used: アクティビティが使用したエンティティ"
    ]
    add_bullet_list(slide, bullets, Inches(0.5), top + Inches(0.3), Inches(9), Inches(1.2))

    add_heading(slide, "生成例", Inches(0.5), top + Inches(1.7), COLOR_TEAL)
    code = """{
  "@context": "http://www.w3.org/ns/prov",
  "entity": {
    "id": "entity:output_001",
    "value": 123.45,
    "unit": "kg-CO2"
  },
  "activity": {
    "id": "activity:energyToEmissions",
    "startedAtTime": "2024-01-15T10:30:00Z"
  },
  "wasGeneratedBy": {
    "entity": "entity:output_001",
    "activity": "activity:energyToEmissions"
  }
}"""
    add_code(slide, code, Inches(0.5), top + Inches(2.0), Inches(9), Inches(1.6))

    return slide

def add_cfp_problem_slide(prs):
    """Slide 18: CFP Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "問題設定")

    top = Inches(1.2)

    add_heading(slide, "シナリオ", Inches(0.5), top, COLOR_TEAL)
    add_text(slide, "企業の購買データから温室効果ガス排出量を計算",
             Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.3))

    add_heading(slide, "入力", Inches(0.5), top + Inches(0.8), COLOR_TEAL)
    bullets1 = ["購買金額（JPY）", "製品カテゴリ", "エネルギー消費量（kWh）"]
    add_bullet_list(slide, bullets1, Inches(0.5), top + Inches(1.1), Inches(9), Inches(0.7))

    add_heading(slide, "目標", Inches(0.5), top + Inches(2.0), COLOR_TEAL)
    bullets2 = ["Scope 1, 2, 3 排出量（kg-CO2）", "総排出量（kg-CO2）"]
    add_bullet_list(slide, bullets2, Inches(0.5), top + Inches(2.3), Inches(9), Inches(0.5))

    add_heading(slide, "課題", Inches(0.5), top + Inches(3.0), COLOR_TEAL)
    bullets3 = [
        "複数の変換経路が存在",
        "集約関数による統合が必要",
        "信頼度とコストのトレードオフ"
    ]
    add_bullet_list(slide, bullets3, Inches(0.5), top + Inches(3.3), Inches(9), Inches(0.8))

    return slide

def add_search_process_slide(prs):
    """Slide 19: Search Process"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "探索プロセス")

    top = Inches(1.2)

    add_heading(slide, "ステップ1：パス発見", Inches(0.5), top, COLOR_TEAL)
    code = """Goal: PurchaseAmount -> TotalGHGEmissions

Found Path:
  PurchaseAmount -> MonetaryValue
  MonetaryValue -> ProductCategory
  ProductCategory -> EnergyConsumption
  EnergyConsumption -> Scope2Emissions
  Scope2Emissions -> TotalGHGEmissions"""
    add_code(slide, code, Inches(0.5), top + Inches(0.3), Inches(9), Inches(1.3))

    add_heading(slide, "ステップ2：最適化", Inches(0.5), top + Inches(1.8), COLOR_TEAL)
    bullets = ["総コスト: 2.1", "総信頼度: 0.85", "中間型数: 4"]
    add_bullet_list(slide, bullets, Inches(0.5), top + Inches(2.1), Inches(9), Inches(0.7))

    add_heading(slide, "ステップ3：実行", Inches(0.5), top + Inches(3.0), COLOR_TEAL)
    add_text(slide, "各関数を順次実行し、Provenanceを記録",
             Inches(0.5), top + Inches(3.3), Inches(9), Inches(0.3))

    return slide

def add_ghg_scope_slide(prs):
    """Slide 20: GHG Scope 123 Scenario"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "GHG Scope 1, 2, 3シナリオ")

    top = Inches(1.2)

    add_heading(slide, "テスト構成", Inches(0.5), top, COLOR_TEAL)
    bullets1 = [
        "Scope 1: 直接排出（燃料燃焼等）",
        "Scope 2: 間接排出（購入電力等）",
        "Scope 3: その他間接排出（サプライチェーン等）"
    ]
    add_bullet_list(slide, bullets1, Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.9))

    add_heading(slide, "8つのテストシナリオ", Inches(0.5), top + Inches(1.4), COLOR_TEAL)
    code = """1. Scope1Emissions → TotalGHGEmissions
2. Scope2Emissions → TotalGHGEmissions
3. Scope3Emissions → TotalGHGEmissions
4. DirectFuelCombustion → Scope1Emissions
5. PurchasedElectricity → Scope2Emissions
6. SupplyChainEmissions → Scope3Emissions
7. DirectFuelCombustion → TotalGHGEmissions (multi-hop)
8. PurchasedElectricity → TotalGHGEmissions (multi-hop)"""
    add_code(slide, code, Inches(0.5), top + Inches(1.7), Inches(9), Inches(1.8))

    return slide

def add_analysis_results_slide(prs):
    """Slide 21: Analysis Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "分析結果")

    top = Inches(1.2)

    add_heading(slide, "成功率", Inches(0.5), top, COLOR_TEAL)
    bullets1 = [
        "100% のシナリオでパス発見成功",
        "平均パス長: 1.5 関数",
        "平均コスト: 1.0"
    ]
    add_bullet_list(slide, bullets1, Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.9))

    add_heading(slide, "集約関数の効果", Inches(0.5), top + Inches(1.4), COLOR_TEAL)
    bullets2 = [
        "複数のScope排出量を統合可能",
        "型システムにシームレスに統合",
        "Provenance記録により透明性確保"
    ]
    add_bullet_list(slide, bullets2, Inches(0.5), top + Inches(1.7), Inches(9), Inches(0.9))

    add_heading(slide, "課題と今後", Inches(0.5), top + Inches(2.8), COLOR_TEAL)
    bullets3 = [
        "現在の実装: 集約関数を単一入力として扱う",
        "理想: 複数引数を持つ依存型システム",
        "拡張案: Product型、Multi-argument関数サポート"
    ]
    add_bullet_list(slide, bullets3, Inches(0.5), top + Inches(3.1), Inches(9), Inches(0.9))

    return slide

def add_metrics_analysis_slide(prs):
    """Slide 22: Metrics Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "メトリクス分析")

    top = Inches(1.2)

    add_heading(slide, "パス長の分布", Inches(0.5), top, COLOR_TEAL)
    bullets1 = [
        "1関数: 6シナリオ (75%)",
        "2関数: 2シナリオ (25%)",
        "最大パス長: 2"
    ]
    add_bullet_list(slide, bullets1, Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.9))

    add_heading(slide, "コストと信頼度", Inches(0.5), top + Inches(1.4), COLOR_TEAL)
    bullets2 = [
        "最小コスト: 1.0",
        "最大コスト: 2.0",
        "平均信頼度: 1.0 (100%)"
    ]
    add_bullet_list(slide, bullets2, Inches(0.5), top + Inches(1.7), Inches(9), Inches(0.9))

    add_heading(slide, "発見", Inches(0.5), top + Inches(2.8), COLOR_TEAL)
    bullets3 = [
        "集約関数はコスト1.0で効率的",
        "マルチホップパスも自動発見可能",
        "型システムが適切に機能している証左"
    ]
    add_bullet_list(slide, bullets3, Inches(0.5), top + Inches(3.1), Inches(9), Inches(0.9))

    return slide

def add_achievements_slide(prs):
    """Slide 24: Achievements"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "主要な成果")

    top = Inches(1.2)

    add_heading(slide, "理論的貢献", Inches(0.5), top, COLOR_TEAL)
    bullets1 = [
        "Span/Cospan理論によるオントロジーアライメントの形式化",
        "型充足問題としての変換経路探索の定式化",
        "中間型の自動発見手法"
    ]
    add_bullet_list(slide, bullets1, Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.9))

    add_heading(slide, "実装的貢献", Inches(0.5), top + Inches(1.4), COLOR_TEAL)
    bullets2 = [
        "YAMLベースのプロトタイプ実装",
        "本格的なDSL設計と実装",
        "実行エンジン（Formula/SPARQL/REST対応）",
        "PROV-O準拠のProvenance生成"
    ]
    add_bullet_list(slide, bullets2, Inches(0.5), top + Inches(1.7), Inches(9), Inches(1.2))

    add_heading(slide, "実証的貢献", Inches(0.5), top + Inches(3.1), COLOR_TEAL)
    bullets3 = [
        "CFP計算での検証",
        "GHG Scope 1,2,3シナリオでの100%成功率",
        "集約関数の統合"
    ]
    add_bullet_list(slide, bullets3, Inches(0.5), top + Inches(3.4), Inches(9), Inches(0.8))

    return slide

def add_future_work_slide(prs):
    """Slide 25: Future Work"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    add_title(slide, "今後の展望")

    top = Inches(1.2)

    add_heading(slide, "型システムの拡張", Inches(0.5), top, COLOR_TEAL)
    bullets1 = [
        "依存型: 値に依存する型（例：長さnのリスト）",
        "多引数関数: (A, B) → C のような関数",
        "Product型: A × B による直積型"
    ]
    add_bullet_list(slide, bullets1, Inches(0.5), top + Inches(0.3), Inches(9), Inches(0.9))

    add_heading(slide, "実行エンジンの強化", Inches(0.5), top + Inches(1.4), COLOR_TEAL)
    bullets2 = [
        "並列実行による高速化",
        "キャッシュ機構の導入",
        "エラーハンドリングの充実",
        "GraphQL/gRPC対応"
    ]
    add_bullet_list(slide, bullets2, Inches(0.5), top + Inches(1.7), Inches(9), Inches(1.2))

    add_heading(slide, "応用領域の拡大", Inches(0.5), top + Inches(3.1), COLOR_TEAL)
    bullets3 = [
        "サプライチェーン全体のCFP計算",
        "生物多様性データの統合",
        "ESG指標の自動計算"
    ]
    add_bullet_list(slide, bullets3, Inches(0.5), top + Inches(3.4), Inches(9), Inches(0.8))

    return slide

def add_conclusion_slide(prs):
    """Slide 26: Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Navy background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "結論"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_LIGHT

    # Main points
    bullets = [
        "型理論とカテゴリ理論により、オントロジー変換を自動化",
        "DSLによる宣言的で拡張可能な設計",
        "実行エンジンにより実用的なシステムを実現",
        "CFP計算での実証により有効性を確認"
    ]

    y_pos = Inches(1.5)
    for bullet in bullets:
        text_box = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(8.4), Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = "• " + bullet
        p = text_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_LIGHT
        p.line_spacing = 1.8
        y_pos += Inches(0.6)

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(0.8))
    footer_frame = footer_box.text_frame
    footer_frame.text = "持続可能性データの統合と分析を\n次のレベルへ"
    footer_p = footer_frame.paragraphs[0]
    footer_p.font.size = Pt(20)
    footer_p.font.color.rgb = COLOR_TEAL
    footer_p.alignment = PP_ALIGN.CENTER

    return slide

# Helper functions

def add_title(slide, title_text):
    """Add title to slide"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_NAVY

def add_heading(slide, text, left, top, color):
    """Add heading text"""
    box = slide.shapes.add_textbox(left, top, Inches(9), Inches(0.4))
    frame = box.text_frame
    frame.text = text
    p = frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color

def add_text(slide, text, left, top, width, height):
    """Add regular text"""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.text = text
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_NAVY

def add_code(slide, code_text, left, top, width, height):
    """Add code block"""
    box = slide.shapes.add_textbox(left, top, width, height)

    # Add background fill
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_CODE_BG

    # Add border
    line = box.line
    line.color.rgb = COLOR_SILVER
    line.width = Pt(1)

    frame = box.text_frame
    frame.text = code_text
    frame.word_wrap = False
    frame.margin_left = Pt(10)
    frame.margin_top = Pt(10)
    frame.margin_right = Pt(10)
    frame.margin_bottom = Pt(10)

    p = frame.paragraphs[0]
    p.font.name = 'Courier New'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_NAVY

def add_bullet_list(slide, items, left, top, width, height):
    """Add bullet list"""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True

    for i, item in enumerate(items):
        if i > 0:
            frame.add_paragraph()
        p = frame.paragraphs[i]
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_NAVY
        p.space_before = Pt(6)
        p.line_spacing = 1.4

if __name__ == '__main__':
    print("Generating presentation...")
    prs = create_presentation()
    output_file = 'slides.pptx'
    prs.save(output_file)
    print(f"✓ Presentation saved to {output_file}")
    print(f"  Total slides: {len(prs.slides)}")
