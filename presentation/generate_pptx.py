#!/usr/bin/env python3
# generate_pptx.py
"""
型理論ベース オントロジー探索・合成システム
プレゼンテーション生成スクリプト (PPTX)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # スライドレイアウト
    title_slide_layout = prs.slide_layouts[0]  # Title Slide
    title_content_layout = prs.slide_layouts[1]  # Title and Content
    blank_layout = prs.slide_layouts[6]  # Blank

    # ========== タイトルスライド ==========
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "型理論ベース\nオントロジー探索・合成システム"
    subtitle.text = "Type-Theoretic Ontology Synthesis System\n\n型理論と圏論に基づくオントロジー変換の自動探索・合成"

    # ========== 目次 ==========
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "目次"

    content = slide.placeholders[1].text_frame
    content.text = "1. 理論的基礎"
    p = content.add_paragraph()
    p.text = "• Span/Cospan理論"
    p.level = 1
    p = content.add_paragraph()
    p.text = "• 型理論・型充足問題"
    p.level = 1
    p = content.add_paragraph()
    p.text = "• 中間型の自動発見"
    p.level = 1

    content.add_paragraph().text = "2. YAMLベースの簡単な実装"
    content.add_paragraph().text = "3. 本格的なDSL設計"
    content.add_paragraph().text = "4. CFP計算の実例"
    content.add_paragraph().text = "5. まとめと今後の展望"

    # ========== セクション1: 理論部分 ==========
    # セクションタイトル
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "1. 理論的基礎"
    content = slide.placeholders[1].text_frame
    content.text = "型理論と圏論によるオントロジー変換の形式化"

    # Span/Cospan理論
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "1.1 Span/Cospan理論"

    content = slide.placeholders[1].text_frame
    content.text = "なぜオントロジーにSpanが向いているのか？"

    p = content.add_paragraph()
    p.text = "オントロジー間のアライメントは本質的に「対応表」を持つ"
    p.level = 1

    p = content.add_paragraph()
    p.text = "→ これがSpanそのもの"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "A ← X → B"
    p.level = 1
    p.font.name = 'Courier New'

    p = content.add_paragraph()
    p.text = "A: オントロジーA, B: オントロジーB"
    p.level = 1
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "X: 対応（マッピングテーブル）"
    p.level = 1
    p.font.size = Pt(14)

    # Spanの合成
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "Spanの合成"

    content = slide.placeholders[1].text_frame
    content.text = "2つのアライメントを連鎖させる"

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "A ← X → B  (AとBのアライメント)"
    p.level = 1
    p.font.name = 'Courier New'

    p = content.add_paragraph()
    p.text = "B ← Y → C  (BとCのアライメント)"
    p.level = 1
    p.font.name = 'Courier New'

    p = content.add_paragraph()
    p.text = "↓ pullback で合成"
    p.level = 1

    p = content.add_paragraph()
    p.text = "A ← X×ᴮY → C  (AとCのアライメント)"
    p.level = 1
    p.font.name = 'Courier New'

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Pullbackの意義:"
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• B型が一致する要素のみが自動的に接合"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• 互換性のないマッピングは合成されない（安全）"
    p.level = 1

    # 型理論
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "1.2 型理論による定式化"

    content = slide.placeholders[1].text_frame
    content.text = "型 = オントロジーのクラス"
    content.paragraphs[0].font.bold = True

    p = content.add_paragraph()
    p.text = "type Product"
    p.level = 1
    p.font.name = 'Courier New'

    p = content.add_paragraph()
    p.text = "type Energy [unit=J]"
    p.level = 1
    p.font.name = 'Courier New'

    p = content.add_paragraph()
    p.text = "type Fuel [unit=kg]"
    p.level = 1
    p.font.name = 'Courier New'

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "関数 = 型間の変換"
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "usesEnergy : Product → Energy"
    p.level = 1
    p.font.name = 'Courier New'

    p = content.add_paragraph()
    p.text = "fuelToCO2 : Fuel → CO2"
    p.level = 1
    p.font.name = 'Courier New'

    # 型充足問題
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "型充足（Type Inhabitation）問題"

    content = slide.placeholders[1].text_frame
    content.text = "問題: 目的型 Product → CO2 を満たす証明項（関数合成）は存在するか？"
    content.paragraphs[0].font.bold = True

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "型付けルールによる導出:"

    p = content.add_paragraph()
    p.text = "Γ ⊢ usesEnergy : Product → Energy"
    p.level = 1
    p.font.name = 'Courier New'
    p.font.size = Pt(12)

    p = content.add_paragraph()
    p.text = "Γ ⊢ energyToFuelEstimate : Energy → Fuel"
    p.level = 1
    p.font.name = 'Courier New'
    p.font.size = Pt(12)

    p = content.add_paragraph()
    p.text = "Γ ⊢ fuelToCO2 : Fuel → CO2"
    p.level = 1
    p.font.name = 'Courier New'
    p.font.size = Pt(12)

    p = content.add_paragraph()
    p.text = "────────────────────────────────"
    p.level = 1
    p.font.name = 'Courier New'

    p = content.add_paragraph()
    p.text = "Γ ⊢ (fuelToCO2 ∘ energyToFuelEstimate ∘ usesEnergy)"
    p.level = 1
    p.font.name = 'Courier New'
    p.font.size = Pt(11)

    p = content.add_paragraph()
    p.text = "   : Product → CO2"
    p.level = 1
    p.font.name = 'Courier New'
    p.font.size = Pt(12)

    # 中間型の自動発見
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "1.3 中間型の自動発見"

    content = slide.placeholders[1].text_frame
    content.text = "探索の過程で「隠れたブリッジ概念」が出現"

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "初期状態: Product と CO2 の直接的な関係は不明"
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "探索プロセス:"
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "1. CO2 を生成する関数は？"
    p.level = 1

    p = content.add_paragraph()
    p.text = "→ fuelToCO2 発見 → Fuel が中間型として出現"
    p.level = 2
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "2. Fuel を生成する関数は？"
    p.level = 1

    p = content.add_paragraph()
    p.text = "→ energyToFuelEstimate 発見 → Energy が出現"
    p.level = 2
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "3. Product から Energy への変換は？"
    p.level = 1

    p = content.add_paragraph()
    p.text = "→ usesEnergy 発見 → 完全なパスが構築"
    p.level = 2
    p.font.bold = True

    # ========== セクション2: YAML実装 ==========
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "2. YAMLベースの実装"
    content = slide.placeholders[1].text_frame
    content.text = "プロトタイプ実装による概念実証"

    # カタログ定義
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "カタログ定義（catalog.yaml）"

    content = slide.placeholders[1].text_frame
    content.text = "types:"
    content.paragraphs[0].font.name = 'Courier New'
    content.paragraphs[0].font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "  - name: Product"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "  - name: Energy"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "    unit: J"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "functions:"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = '  - id: usesEnergy'
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = '    sig: "Product -> Energy"'
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = '    cost: 1'
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = '    confidence: 0.9'
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    # 探索アルゴリズム
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "探索アルゴリズム: 逆方向探索"

    content = slide.placeholders[1].text_frame
    content.text = "アルゴリズム概要:"
    content.paragraphs[0].font.bold = True

    p = content.add_paragraph()
    p.text = "1. ゴール型からスタート"
    p.level = 1

    p = content.add_paragraph()
    p.text = "2. ゴールを返す関数を探索"
    p.level = 1

    p = content.add_paragraph()
    p.text = "3. 各関数のdomainを新しいサブゴールに"
    p.level = 1

    p = content.add_paragraph()
    p.text = "4. ソース型に到達したらパス記録"
    p.level = 1

    p = content.add_paragraph()
    p.text = "5. コスト最小のパスを返す"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "特徴:"
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• Dijkstra的な最短経路探索"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• 優先度付きキュー使用"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• コスト制限による枝刈り"
    p.level = 1

    # 検証結果
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "検証結果"

    content = slide.placeholders[1].text_frame
    content.text = "すべてのテストが成功"
    content.paragraphs[0].font.bold = True
    content.paragraphs[0].font.size = Pt(24)
    content.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "✓ 基本動作: Product → CO2 のパス探索成功"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ 包括的テスト: 6つのテストケースすべて成功"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ DSL統合テスト: 5つのテストケースすべて成功"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ 実行機能テスト: 5つのテストケースすべて成功"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ GHG集約テスト: 8つのシナリオすべて成功"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ 型理論・圏論との整合性確認"
    p.level = 1

    # ========== セクション3: 本格的なDSL ==========
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "3. 本格的なDSL設計"
    content = slide.placeholders[1].text_frame
    content.text = "読みやすく拡張可能な専用言語"

    # DSL構文
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "DSL構文"

    content = slide.placeholders[1].text_frame
    content.text = "# 型定義"
    content.paragraphs[0].font.name = 'Courier New'
    content.paragraphs[0].font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "type Product"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "type Energy [unit=J, range=>=0]"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "# 関数定義"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "fn usesEnergy {"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "  sig: Product -> Energy"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = '  impl: sparql("SELECT ...")'
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "  cost: 1"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "  confidence: 0.9"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "}"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    # 実行機能
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "実行機能（新実装）"

    content = slide.placeholders[1].text_frame
    content.text = "3つの主要機能を追加:"
    content.paragraphs[0].font.bold = True

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "1. 実行レイヤー: SPARQL/REST/Formula の実際の実行"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• FormulaExecutor: 数式の安全な評価"
    p.level = 2

    p = content.add_paragraph()
    p.text = "• SPARQLExecutor: SPARQLクエリ実行"
    p.level = 2

    p = content.add_paragraph()
    p.text = "• RESTExecutor: REST API呼び出し"
    p.level = 2

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "2. 単位変換: 自動的な単位変換関数の挿入"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• J ↔ kWh, kg ↔ g, K ↔ C ↔ F など"
    p.level = 2

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "3. Provenance: PROV-O形式での来歴記録"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Turtle および JSON 形式で出力"
    p.level = 2

    # ========== セクション4: CFP事例 ==========
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "4. CFP計算の実例"
    content = slide.placeholders[1].text_frame
    content.text = "Carbon Footprint 計算による実証"

    # 問題設定
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "問題設定"

    content = slide.placeholders[1].text_frame
    content.text = "目標: 製品（Product）のCO2排出量を計算"
    content.paragraphs[0].font.bold = True

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "制約: 直接的な Product → CO2 関数は存在しない"
    p.font.bold = True

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "利用可能な関数（断片的）:"

    p = content.add_paragraph()
    p.text = "• usesEnergy: Product → Energy"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• fuelToEnergy: Fuel → Energy"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• fuelToCO2: Fuel → CO2"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• energyToFuelEstimate: Energy → Fuel (逆関数)"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "→ これらをどう組み合わせるか？"
    p.font.italic = True

    # 発見されたパス
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "発見されたパス"

    content = slide.placeholders[1].text_frame
    content.text = "Product"
    content.paragraphs[0].font.name = 'Courier New'
    content.paragraphs[0].font.size = Pt(16)

    p = content.add_paragraph()
    p.text = "   |"
    p.font.name = 'Courier New'
    p.font.size = Pt(16)

    p = content.add_paragraph()
    p.text = "   | [1] usesEnergy (conf: 0.9)"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "   ↓"
    p.font.name = 'Courier New'
    p.font.size = Pt(16)

    p = content.add_paragraph()
    p.text = "Energy"
    p.font.name = 'Courier New'
    p.font.size = Pt(16)

    p = content.add_paragraph()
    p.text = "   |"
    p.font.name = 'Courier New'
    p.font.size = Pt(16)

    p = content.add_paragraph()
    p.text = "   | [3] energyToFuelEstimate (conf: 0.8) ⚠️"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "   ↓"
    p.font.name = 'Courier New'
    p.font.size = Pt(16)

    p = content.add_paragraph()
    p.text = "Fuel"
    p.font.name = 'Courier New'
    p.font.size = Pt(16)

    p = content.add_paragraph()
    p.text = "   |"
    p.font.name = 'Courier New'
    p.font.size = Pt(16)

    p = content.add_paragraph()
    p.text = "   | [1] fuelToCO2 (conf: 0.98)"
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "   ↓"
    p.font.name = 'Courier New'
    p.font.size = Pt(16)

    p = content.add_paragraph()
    p.text = "CO2"
    p.font.name = 'Courier New'
    p.font.size = Pt(16)
    p.font.bold = True

    # メトリクス分析
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "メトリクス分析"

    content = slide.placeholders[1].text_frame
    content.text = "総コスト: 5.0"
    content.paragraphs[0].font.bold = True
    content.paragraphs[0].font.size = Pt(20)

    p = content.add_paragraph()
    p.text = "• usesEnergy: 1.0 (20%)"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• energyToFuelEstimate: 3.0 (60%) ← 支配的"
    p.level = 1
    p.font.color.rgb = RGBColor(255, 100, 0)

    p = content.add_paragraph()
    p.text = "• fuelToCO2: 1.0 (20%)"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "総信頼度: 0.7056 ≈ 70.56%"
    p.font.bold = True
    p.font.size = Pt(20)

    p = content.add_paragraph()
    p.text = "= 0.9 × 0.8 × 0.98"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "→ 実用的に許容範囲"
    p.font.italic = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    # GHG Scope
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "GHG Scope 1, 2, 3 シナリオ"

    content = slide.placeholders[1].text_frame
    content.text = "集約関数の影響分析"
    content.paragraphs[0].font.bold = True

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Scope 1: 直接排出（燃料燃焼、プロセス排出）"
    p.level = 1

    p = content.add_paragraph()
    p.text = "Scope 2: エネルギー間接排出（購入電力、購入熱）"
    p.level = 1

    p = content.add_paragraph()
    p.text = "Scope 3: その他間接排出（輸送、通勤、出張）"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "テスト結果: 8つのシナリオすべて成功"
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "発見: 集約関数は通常の関数と同様に型合成に統合され、"
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "パス探索に影響を与えない（成功率100%）"
    p.font.size = Pt(14)

    # ========== セクション5: まとめ ==========
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "5. まとめと今後の展望"
    content = slide.placeholders[1].text_frame
    content.text = ""

    # 主要な成果
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "主要な成果"

    content = slide.placeholders[1].text_frame
    content.text = "理論と実装の統合"
    content.paragraphs[0].font.bold = True
    content.paragraphs[0].font.size = Pt(20)

    p = content.add_paragraph()
    p.text = "✓ Span/Cospan理論の実装可能性を実証"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ 型充足問題としての定式化が有効"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ 中間型の自動発見メカニズムを確認"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "実用的なシステム"
    p.font.bold = True
    p.font.size = Pt(20)

    p = content.add_paragraph()
    p.text = "✓ 自動パス探索（コスト最適化）"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ 専用DSL（読みやすく拡張可能）"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ 実行レイヤー（SPARQL/REST/Formula）"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ 単位変換・Provenance生成"
    p.level = 1

    # 実装済み機能
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "実装済み機能"

    content = slide.placeholders[1].text_frame
    content.text = "コア機能:"
    content.paragraphs[0].font.bold = True

    p = content.add_paragraph()
    p.text = "• 型合成エンジン"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• 逆方向探索アルゴリズム"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• コスト・信頼度最適化"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "拡張機能:"
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• 実行レイヤー"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• 単位変換システム"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Provenance生成"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• モックモード（テスト用）"
    p.level = 1

    # 今後の拡張
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "今後の拡張"

    content = slide.placeholders[1].text_frame
    content.text = "中期: アルゴリズム改善"
    content.paragraphs[0].font.bold = True

    p = content.add_paragraph()
    p.text = "• A*探索: ヒューリスティクス関数の導入"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• 双方向探索: 起点と終点からの同時探索"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• キャッシング: 型到達可能性のキャッシュ"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "長期: 理論的拡張"
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• 依存型: 値に依存する型"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• 多引数関数: カリー化を超えた表現力"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Cospan/Colimit: オントロジー統合"
    p.level = 1

    # まとめ
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "まとめ"

    content = slide.placeholders[1].text_frame
    content.text = "型理論と圏論を実装レベルで統合"
    content.paragraphs[0].font.bold = True
    content.paragraphs[0].font.size = Pt(20)

    p = content.add_paragraph()
    p.text = "抽象的な理論を具体的な問題解決に応用"
    p.font.italic = True

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "✓ 理論的健全性と実用性を両立"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ 自動探索による生産性向上"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ 完全なトレーサビリティ"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ 拡張可能なアーキテクチャ"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "GitHub: github.com/miwamasa/"
    p.font.size = Pt(14)

    p = content.add_paragraph()
    p.text = "        ccw-type_theoretic_ontology_part2"
    p.font.size = Pt(14)

    # 最終スライド
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "ありがとうございました"
    subtitle.text = "Type-Theoretic Ontology Synthesis System\n型理論ベース オントロジー探索・合成システム\n\n質問・フィードバックをお待ちしています"

    return prs

def main():
    print("Generating PPTX presentation...")
    prs = create_presentation()

    output_path = '/home/user/ccw-type_theoretic_ontology_part2/presentation/slides.pptx'
    prs.save(output_path)
    print(f"✓ PPTX saved to: {output_path}")

    # カウント
    print(f"✓ Total slides: {len(prs.slides)}")

if __name__ == '__main__':
    main()
